"""MD xtract: universal extractor to clean markdown from files/URLs."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, Optional, Tuple
from urllib.parse import urlparse

import httpx


_CONTROL_CHARS = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]")
_DANGEROUS_TOKENS = [
    "<system>",
    "</system>",
    "<assistant>",
    "</assistant>",
    "<user>",
    "</user>",
]


def _sanitize_markdown(value: str, max_chars: int = 0) -> str:
    """Sanitize text but preserve newlines for markdown."""
    if value is None:
        return ""
    cleaned = _CONTROL_CHARS.sub("", str(value))
    for token in _DANGEROUS_TOKENS:
        cleaned = cleaned.replace(token, f"[{token.strip('<>')}]")
    cleaned = cleaned.replace("```", "'''")
    if max_chars and len(cleaned) > max_chars:
        cleaned = cleaned[: max_chars - 1] + "…"
    return cleaned


def _is_url(value: str) -> bool:
    try:
        parsed = urlparse(value)
        return parsed.scheme in ("http", "https")
    except Exception:
        return False


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf_bytes(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as e:
        raise RuntimeError("pypdf is required to parse PDFs. Install with: pip install pypdf") from e
    from io import BytesIO
    reader = PdfReader(BytesIO(data))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


def _read_pdf(path: Path) -> str:
    return _read_pdf_bytes(path.read_bytes())


def _read_docx(path: Path) -> str:
    try:
        import docx2txt
    except Exception as e:
        raise RuntimeError("docx2txt is required for DOCX. Install with: pip install docx2txt") from e
    return docx2txt.process(str(path)) or ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("python-pptx is required for PPTX. Install with: pip install python-pptx") from e
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _read_xlsx(path: Path) -> str:
    try:
        import openpyxl
    except Exception as e:
        raise RuntimeError("openpyxl is required for XLSX. Install with: pip install openpyxl") from e
    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_vals = ["" if v is None else str(v) for v in row]
            rows.append(row_vals)
        if not rows:
            continue
        header = rows[0]
        parts.append(f"### Sheet: {sheet_name}")
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for r in rows[1:]:
            parts.append("| " + " | ".join(r) + " |")
        parts.append("")
    return "\n".join(parts)


def _read_html_bytes(data: bytes) -> Tuple[str, Dict[str, Any]]:
    text = data.decode("utf-8", errors="ignore")
    title = None
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else None
        body = str(soup.body or soup)
    except Exception:
        body = text
    try:
        import markdownify  # type: ignore
        md = markdownify.markdownify(body, heading_style="ATX")
    except Exception:
        md = re.sub(r"<[^>]+>", " ", body)
    return md, {"title": title}


def _read_html(path: Path) -> Tuple[str, Dict[str, Any]]:
    return _read_html_bytes(path.read_bytes())


def _read_image(path: Path, ocr: bool) -> str:
    if not ocr:
        raise RuntimeError("OCR disabled. Pass ocr=true to extract text from images.")
    try:
        from PIL import Image
    except Exception as e:
        raise RuntimeError("Pillow is required for images. Install with: pip install pillow") from e
    try:
        import pytesseract
    except Exception as e:
        raise RuntimeError("pytesseract is required for OCR. Install with: pip install pytesseract") from e
    image = Image.open(str(path))
    return pytesseract.image_to_string(image)


def xtract_to_markdown(
    input_ref: str,
    input_type: Optional[str] = None,
    max_chars: int = 0,
    ocr: bool = False,
) -> Tuple[str, Dict[str, Any]]:
    """Extract clean markdown from a file path or URL."""
    meta: Dict[str, Any] = {"input": input_ref}
    source_type = input_type or ("url" if _is_url(input_ref) else "path")
    meta["source_type"] = source_type

    if source_type == "url":
        with httpx.Client(timeout=30, follow_redirects=True) as client:
            resp = client.get(input_ref)
            resp.raise_for_status()
            content_type = (resp.headers.get("content-type") or "").lower()
            data = resp.content
        meta["content_type"] = content_type
        if "application/pdf" in content_type or input_ref.lower().endswith(".pdf"):
            markdown = _read_pdf_bytes(data)
        elif "text/html" in content_type or input_ref.lower().endswith((".html", ".htm")):
            markdown, html_meta = _read_html_bytes(data)
            meta.update(html_meta)
        else:
            markdown = data.decode("utf-8", errors="ignore")
    else:
        path = Path(input_ref).expanduser().resolve()
        if not path.exists():
            raise FileNotFoundError(f"Path not found: {path}")
        meta["file_path"] = str(path)
        suffix = path.suffix.lower()
        if suffix in (".md", ".txt"):
            markdown = _read_text_file(path)
        elif suffix == ".pdf":
            markdown = _read_pdf(path)
        elif suffix == ".docx":
            markdown = _read_docx(path)
        elif suffix == ".pptx":
            markdown = _read_pptx(path)
        elif suffix == ".xlsx":
            markdown = _read_xlsx(path)
        elif suffix in (".html", ".htm"):
            markdown, html_meta = _read_html(path)
            meta.update(html_meta)
        elif suffix in (".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"):
            markdown = _read_image(path, ocr=ocr)
        else:
            markdown = _read_text_file(path)

    markdown = _sanitize_markdown(markdown, max_chars=max_chars)
    meta["chars"] = len(markdown)
    return markdown, meta
